"""Tests for PptxMerger — generates test fixtures, no external files needed."""

import io
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from pptx_slide_merger import PptxMerger, list_slides


def _make_pptx(slide_titles: list[str], tmp_dir: Path) -> Path:
    """Create a simple PPTX with titled slides."""
    prs = PptxPresentation()
    for title_text in slide_titles:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = title_text
    path = tmp_dir / f"test_{'_'.join(t[:8] for t in slide_titles)}.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def tmp_dir():
    with tempfile.TemporaryDirectory() as d:
        yield Path(d)


class TestListSlides:
    def test_list_slides_returns_titles(self, tmp_dir):
        pptx = _make_pptx(["Alpha", "Beta", "Gamma"], tmp_dir)
        titles = list_slides(pptx)
        assert len(titles) == 3
        assert titles[0] == "Alpha"
        assert titles[1] == "Beta"
        assert titles[2] == "Gamma"

    def test_list_slides_empty_deck(self, tmp_dir):
        prs = PptxPresentation()
        path = tmp_dir / "empty.pptx"
        prs.save(str(path))
        titles = list_slides(path)
        assert titles == []


class TestPptxMerger:
    def test_merge_single_source(self, tmp_dir):
        base = _make_pptx(["Base1", "Base2"], tmp_dir)
        output = tmp_dir / "merged.pptx"

        with PptxMerger(base, verbose=False) as merger:
            assert merger.add_slide(base, 0) is True
            assert merger.add_slide(base, 1) is True
            merger.save(output)

        result = PptxPresentation(str(output))
        assert len(result.slides) == 2

    def test_merge_multiple_sources(self, tmp_dir):
        deck_a = _make_pptx(["A1", "A2", "A3"], tmp_dir)
        deck_b = _make_pptx(["B1", "B2"], tmp_dir)
        output = tmp_dir / "merged.pptx"

        with PptxMerger(deck_a, verbose=False) as merger:
            merger.add_slide(deck_a, 0)
            merger.add_slide(deck_b, 0)
            merger.add_slide(deck_a, 2)
            merger.add_slide(deck_b, 1)
            merger.save(output)

        result = PptxPresentation(str(output))
        assert len(result.slides) == 4

    def test_invalid_slide_index_returns_false(self, tmp_dir):
        base = _make_pptx(["Only"], tmp_dir)
        with PptxMerger(base, verbose=False) as merger:
            assert merger.add_slide(base, 99) is False

    def test_output_is_valid_pptx(self, tmp_dir):
        base = _make_pptx(["Slide1", "Slide2", "Slide3"], tmp_dir)
        output = tmp_dir / "valid.pptx"

        with PptxMerger(base, verbose=False) as merger:
            merger.add_slide(base, 1)
            merger.save(output)

        # python-pptx should open it without error
        result = PptxPresentation(str(output))
        assert len(result.slides) == 1

    def test_merge_preserves_slide_content(self, tmp_dir):
        base = _make_pptx(["Hello World"], tmp_dir)
        output = tmp_dir / "content.pptx"

        with PptxMerger(base, verbose=False) as merger:
            merger.add_slide(base, 0)
            merger.save(output)

        result = PptxPresentation(str(output))
        slide = result.slides[0]
        assert slide.shapes.title is not None
        assert slide.shapes.title.text == "Hello World"

    def test_context_manager_cleanup(self, tmp_dir):
        base = _make_pptx(["Test"], tmp_dir)
        with PptxMerger(base, verbose=False) as merger:
            tmp = merger.tmp_dir
        # After exiting, tmp_dir should be cleaned up
        assert not Path(tmp).exists()

    def test_merge_with_image(self, tmp_dir):
        """Verify slides with images merge without error."""
        prs = PptxPresentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

        # Create a tiny valid PNG (1x1 red pixel)
        import struct
        import zlib

        def _make_png():
            raw = b'\x00\xff\x00\x00'  # filter byte + RGB
            compressed = zlib.compress(raw)
            chunks = b''
            for tag, data in [
                (b'IHDR', struct.pack('>IIBBBBB', 1, 1, 8, 2, 0, 0, 0)),
                (b'IDAT', compressed),
                (b'IEND', b''),
            ]:
                chunk = tag + data
                crc = struct.pack('>I', zlib.crc32(chunk) & 0xffffffff)
                chunks += struct.pack('>I', len(data)) + chunk + crc
            return b'\x89PNG\r\n\x1a\n' + chunks

        png_path = tmp_dir / "test.png"
        png_path.write_bytes(_make_png())
        slide.shapes.add_picture(str(png_path), Inches(1), Inches(1),
                                 Inches(2), Inches(2))

        base_path = tmp_dir / "with_image.pptx"
        prs.save(str(base_path))

        output = tmp_dir / "merged_image.pptx"
        with PptxMerger(base_path, verbose=False) as merger:
            merger.add_slide(base_path, 0)
            merger.save(output)

        result = PptxPresentation(str(output))
        assert len(result.slides) == 1
