"""
Core PPTX slide merger implementation.

Copies slides between PowerPoint files preserving each slide's original
layout, master, and theme. Works at the ZIP/XML level — no .NET required.

Key insight: OOXML uses a global ID space for sldMasterId and sldLayoutId
values. When copying masters from different source decks, these IDs must
be remapped to avoid collisions.
"""

import os
import re
import shutil
import tempfile
import zipfile
from pathlib import Path
from lxml import etree

from pptx import Presentation as PptxPresentation

# OOXML namespace URIs
NS_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# Content types
SLIDE_CT = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
LAYOUT_CT = "application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"
MASTER_CT = "application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"
NOTES_CT = "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"
THEME_CT = "application/vnd.openxmlformats-officedocument.theme+xml"
NOTES_MASTER_CT = "application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml"

# Relationship types
REL_SLIDE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
REL_LAYOUT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
REL_MASTER = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster"
REL_IMAGE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
REL_NOTES = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
REL_THEME = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
REL_NOTES_MASTER = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster"
REL_VIDEO = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/video"
REL_AUDIO = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio"
REL_MEDIA_ALT = "http://schemas.microsoft.com/office/2007/relationships/media"

MEDIA_REL_TYPES = {REL_IMAGE, REL_VIDEO, REL_AUDIO, REL_MEDIA_ALT}
MEDIA_REL_KEYWORDS = {'image', 'video', 'audio', 'media'}

MEDIA_MIME = {
    'png': 'image/png', 'jpg': 'image/jpeg', 'jpeg': 'image/jpeg',
    'gif': 'image/gif', 'svg': 'image/svg+xml', 'emf': 'image/x-emf',
    'wmf': 'image/x-wmf', 'tiff': 'image/tiff', 'tif': 'image/tiff',
    'wdp': 'image/vnd.ms-photo', 'fntdata': 'application/x-fontdata',
    'mp4': 'video/mp4', 'm4v': 'video/mp4', 'mov': 'video/quicktime',
    'avi': 'video/avi', 'wmv': 'video/x-ms-wmv', 'mp3': 'audio/mpeg',
    'wav': 'audio/wav', 'wma': 'audio/x-ms-wma',
}

# OOXML element ordering for <p:presentation>
PRES_ELEMENT_ORDER = [
    'sldMasterIdLst', 'notesMasterIdLst', 'handoutMasterIdLst',
    'sldIdLst', 'sldSz', 'notesSz', 'smartTags', 'embeddedFontLst',
    'custShowLst', 'photoAlbum', 'custDataLst', 'kinsoku',
    'defaultTextStyle', 'modifyVerifier',
]


def list_slides(pptx_path: Path) -> list[str]:
    """Return a list of slide titles from a PPTX file."""
    prs = PptxPresentation(str(pptx_path))
    titles = []
    for slide in prs.slides:
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.replace("\n", " ").strip()
        if not title:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    title = shape.text_frame.text[:80].replace("\n", " ").strip()
                    break
        titles.append(title or "(no title)")
    return titles


def _parse_xml(path):
    return etree.parse(path)


def _write_xml(tree, path):
    tree.write(path, xml_declaration=True, encoding="UTF-8", standalone=True)


def _resolve(base_dir, relative_target):
    return os.path.normpath(os.path.join(base_dir, relative_target))


def _max_rid(rels_root):
    mx = 0
    for r in rels_root.findall(f'{{{NS_REL}}}Relationship'):
        m = re.match(r'rId(\d+)', r.get('Id', ''))
        if m:
            mx = max(mx, int(m.group(1)))
    return mx


def _is_media_rel(rel_type: str) -> bool:
    """Check if a relationship type refers to a media file (image/video/audio)."""
    if rel_type in MEDIA_REL_TYPES:
        return True
    lower = rel_type.lower().split('/')[-1]
    return lower in MEDIA_REL_KEYWORDS


class PptxMerger:
    """
    Merge slides from multiple PPTX files, preserving each slide's
    original layout, master, and theme.

    Usage::

        merger = PptxMerger(Path("base.pptx"))
        merger.add_slide(Path("source1.pptx"), 0)  # slide index 0
        merger.add_slide(Path("source2.pptx"), 3)  # slide index 3
        merger.add_slide(Path("base.pptx"), 5)      # from base deck too
        merger.save(Path("output.pptx"))
        merger.cleanup()

    The base PPTX provides the slide dimensions and its masters/layouts
    are kept intact. Slides from the base deck reuse existing layouts.
    Slides from other decks get their masters (with all layouts and themes)
    copied in alongside.
    """

    def __init__(self, base_pptx: Path, verbose: bool = True):
        self._verbose = verbose
        self.tmp_dir = tempfile.mkdtemp(prefix="pptx_merge_")
        self.base_dir = os.path.join(self.tmp_dir, "base")
        self._src_cache = {}

        with zipfile.ZipFile(base_pptx, 'r') as z:
            z.extractall(self.base_dir)

        self.pres_xml_path = os.path.join(self.base_dir, "ppt", "presentation.xml")
        self.pres_rels_path = os.path.join(self.base_dir, "ppt", "_rels", "presentation.xml.rels")
        self.ct_path = os.path.join(self.base_dir, "[Content_Types].xml")

        self.pres_tree = _parse_xml(self.pres_xml_path)
        self.rels_tree = _parse_xml(self.pres_rels_path)
        self.ct_tree = _parse_xml(self.ct_path)

        self._strip_slides()

        self._next_slide = 1
        self._next_layout = self._count_files("ppt/slideLayouts", "slideLayout") + 1
        self._next_master = self._count_files("ppt/slideMasters", "slideMaster") + 1
        self._next_theme = self._count_files("ppt/theme", "theme") + 1
        self._next_notes = 1
        self._next_media = self._count_media() + 1
        self._next_slide_id = 256
        self._next_pres_rid = _max_rid(self.rels_tree.getroot()) + 1
        self._next_layout_id = self._find_max_layout_id() + 1

        self._layout_map = {}
        self._master_map = {}
        self._media_hash_map = {}  # content_hash -> relative path from ppt/
        self._index_existing_media()
        self._notes_master_copied = self._has_notes_master()
        self._register_base_layouts(base_pptx)
        self._slide_count = 0

        if self._verbose:
            print(f"Base: {base_pptx.name} (layouts={self._next_layout-1}, "
                  f"masters={self._next_master-1})")

    def _strip_slides(self):
        pres_root = self.pres_tree.getroot()
        rels_root = self.rels_tree.getroot()
        ct_root = self.ct_tree.getroot()

        sldIdLst = pres_root.find(f'{{{NS_P}}}sldIdLst')
        if sldIdLst is not None:
            for child in list(sldIdLst):
                sldIdLst.remove(child)

        for rel in list(rels_root.findall(f'{{{NS_REL}}}Relationship')):
            if rel.get('Type') == REL_SLIDE:
                rels_root.remove(rel)

        for override in list(ct_root.findall(f'{{{NS_CT}}}Override')):
            pn = override.get('PartName', '')
            if '/ppt/slides/slide' in pn or '/ppt/notesSlides/' in pn:
                ct_root.remove(override)

        for subdir in ('slides', 'notesSlides'):
            d = os.path.join(self.base_dir, "ppt", subdir)
            if os.path.exists(d):
                shutil.rmtree(d)

    def _count_files(self, subdir, prefix):
        d = os.path.join(self.base_dir, subdir)
        if not os.path.exists(d):
            return 0
        return len([f for f in os.listdir(d)
                    if f.startswith(prefix) and f.endswith('.xml')])

    def _count_media(self):
        d = os.path.join(self.base_dir, "ppt", "media")
        return len(os.listdir(d)) if os.path.exists(d) else 0

    def _find_max_layout_id(self):
        mx = 2147483648
        masters_dir = os.path.join(self.base_dir, "ppt", "slideMasters")
        if not os.path.exists(masters_dir):
            return mx
        for f in os.listdir(masters_dir):
            if f.startswith("slideMaster") and f.endswith(".xml"):
                tree = etree.parse(os.path.join(masters_dir, f))
                for li in tree.getroot().findall(f'.//{{{NS_P}}}sldLayoutId'):
                    val = int(li.get('id', '0'))
                    if val > mx:
                        mx = val
        return mx

    def _index_existing_media(self):
        """Hash all existing media files in the base so we can deduplicate."""
        import hashlib
        media_dir = os.path.join(self.base_dir, "ppt", "media")
        if not os.path.exists(media_dir):
            return
        for f in os.listdir(media_dir):
            fp = os.path.join(media_dir, f)
            if os.path.isfile(fp):
                h = hashlib.md5(open(fp, 'rb').read()).hexdigest()
                self._media_hash_map[h] = f"media/{f}"

    def _has_notes_master(self):
        return os.path.exists(
            os.path.join(self.base_dir, "ppt", "notesMasters", "notesMaster1.xml"))

    def _register_base_layouts(self, base_pptx: Path):
        self._extract_source(base_pptx)
        layouts_dir = os.path.join(self.base_dir, "ppt", "slideLayouts")
        if os.path.exists(layouts_dir):
            for f in os.listdir(layouts_dir):
                if f.startswith("slideLayout") and f.endswith(".xml"):
                    self._layout_map[(str(base_pptx), f)] = f

        masters_dir = os.path.join(self.base_dir, "ppt", "slideMasters")
        if os.path.exists(masters_dir):
            for f in os.listdir(masters_dir):
                if f.startswith("slideMaster") and f.endswith(".xml"):
                    self._master_map[(str(base_pptx), f)] = f

    def _alloc_pres_rid(self):
        rid = f"rId{self._next_pres_rid}"
        self._next_pres_rid += 1
        return rid

    def _extract_source(self, source_pptx: Path) -> str:
        key = str(source_pptx)
        if key not in self._src_cache:
            d = os.path.join(self.tmp_dir, f"src_{len(self._src_cache)}")
            with zipfile.ZipFile(source_pptx, 'r') as z:
                z.extractall(d)
            self._src_cache[key] = d
        return self._src_cache[key]

    def _copy_media_file(self, src_media_path: str) -> str:
        if not os.path.exists(src_media_path):
            return None
        import hashlib
        content = open(src_media_path, 'rb').read()
        h = hashlib.md5(content).hexdigest()
        if h in self._media_hash_map:
            return self._media_hash_map[h]
        ext = os.path.splitext(src_media_path)[1]
        new_name = f"image{self._next_media}{ext}"
        self._next_media += 1
        dst = os.path.join(self.base_dir, "ppt", "media", new_name)
        os.makedirs(os.path.dirname(dst), exist_ok=True)
        shutil.copy2(src_media_path, dst)
        self._ensure_ext_ct(ext.lstrip('.'))
        self._media_hash_map[h] = f"media/{new_name}"
        return f"media/{new_name}"

    def _ensure_ext_ct(self, ext):
        if ext.lower() not in MEDIA_MIME:
            return
        ct_root = self.ct_tree.getroot()
        for d in ct_root.findall(f'{{{NS_CT}}}Default'):
            if d.get('Extension', '').lower() == ext.lower():
                return
        el = etree.SubElement(ct_root, f'{{{NS_CT}}}Default')
        el.set('Extension', ext.lower())
        el.set('ContentType', MEDIA_MIME[ext.lower()])

    def _add_ct_override(self, part_name, content_type):
        ct_root = self.ct_tree.getroot()
        for o in ct_root.findall(f'{{{NS_CT}}}Override'):
            if o.get('PartName') == part_name:
                return
        el = etree.SubElement(ct_root, f'{{{NS_CT}}}Override')
        el.set('PartName', part_name)
        el.set('ContentType', content_type)

    def _remap_rels_media(self, rels_root, src_part_dir):
        for rel in rels_root.findall(f'{{{NS_REL}}}Relationship'):
            rel_type = rel.get('Type', '')
            if _is_media_rel(rel_type):
                src_path = _resolve(src_part_dir, rel.get('Target', ''))
                new_path = self._copy_media_file(src_path)
                if new_path:
                    rel.set('Target', f"../{new_path}")

    def _copy_notes_master(self, src_dir: str):
        if self._notes_master_copied:
            return
        src_nm = os.path.join(src_dir, "ppt", "notesMasters", "notesMaster1.xml")
        if not os.path.exists(src_nm):
            return
        dst_nm = os.path.join(self.base_dir, "ppt", "notesMasters", "notesMaster1.xml")
        os.makedirs(os.path.dirname(dst_nm), exist_ok=True)
        shutil.copy2(src_nm, dst_nm)

        src_rels = os.path.join(src_dir, "ppt", "notesMasters", "_rels",
                                "notesMaster1.xml.rels")
        if os.path.exists(src_rels):
            rels_tree = _parse_xml(src_rels)
            rels_root = rels_tree.getroot()
            for rel in rels_root.findall(f'{{{NS_REL}}}Relationship'):
                rel_type = rel.get('Type', '')
                target = rel.get('Target', '')
                if _is_media_rel(rel_type):
                    src_path = _resolve(
                        os.path.join(src_dir, "ppt", "notesMasters"), target)
                    new_path = self._copy_media_file(src_path)
                    if new_path:
                        rel.set('Target', f"../{new_path}")
                elif rel_type == REL_THEME:
                    src_theme = _resolve(
                        os.path.join(src_dir, "ppt", "notesMasters"), target)
                    if os.path.exists(src_theme):
                        new_theme = f"theme{self._next_theme}.xml"
                        self._next_theme += 1
                        dst_theme = os.path.join(
                            self.base_dir, "ppt", "theme", new_theme)
                        os.makedirs(os.path.dirname(dst_theme), exist_ok=True)
                        shutil.copy2(src_theme, dst_theme)
                        rel.set('Target', f"../theme/{new_theme}")
                        self._add_ct_override(f"/ppt/theme/{new_theme}", THEME_CT)
            dst_rels = os.path.join(self.base_dir, "ppt", "notesMasters",
                                    "_rels", "notesMaster1.xml.rels")
            os.makedirs(os.path.dirname(dst_rels), exist_ok=True)
            _write_xml(rels_tree, dst_rels)

        self._add_ct_override("/ppt/notesMasters/notesMaster1.xml", NOTES_MASTER_CT)

        rid = self._alloc_pres_rid()
        pres_rels_root = self.rels_tree.getroot()
        rel_el = etree.SubElement(pres_rels_root, f'{{{NS_REL}}}Relationship')
        rel_el.set('Id', rid)
        rel_el.set('Type', REL_NOTES_MASTER)
        rel_el.set('Target', "notesMasters/notesMaster1.xml")

        pres_root = self.pres_tree.getroot()
        nmIdLst = pres_root.find(f'{{{NS_P}}}notesMasterIdLst')
        if nmIdLst is None:
            nmIdLst = etree.SubElement(pres_root, f'{{{NS_P}}}notesMasterIdLst')
        nm_el = etree.SubElement(nmIdLst, f'{{{NS_P}}}notesMasterId')
        nm_el.set(f'{{{NS_R}}}id', rid)

        self._notes_master_copied = True

    def _copy_master_with_all_layouts(self, source_pptx: Path, src_dir: str,
                                       master_filename: str) -> str:
        key = (str(source_pptx), master_filename)
        if key in self._master_map:
            return self._master_map[key]

        new_master = f"slideMaster{self._next_master}.xml"
        self._next_master += 1

        src_master_path = os.path.join(src_dir, "ppt", "slideMasters", master_filename)
        dst_master_path = os.path.join(self.base_dir, "ppt", "slideMasters", new_master)
        os.makedirs(os.path.dirname(dst_master_path), exist_ok=True)
        shutil.copy2(src_master_path, dst_master_path)

        master_xml_tree = _parse_xml(dst_master_path)
        master_xml_root = master_xml_tree.getroot()

        src_rels_path = os.path.join(src_dir, "ppt", "slideMasters", "_rels",
                                     f"{master_filename}.rels")
        dst_rels_path = os.path.join(self.base_dir, "ppt", "slideMasters", "_rels",
                                     f"{new_master}.rels")

        layout_rid_map = {}

        if os.path.exists(src_rels_path):
            rels_tree = _parse_xml(src_rels_path)
            rels_root = rels_tree.getroot()

            for rel in rels_root.findall(f'{{{NS_REL}}}Relationship'):
                target = rel.get('Target', '')
                rel_type = rel.get('Type', '')
                rid = rel.get('Id', '')

                if _is_media_rel(rel_type):
                    src_path = _resolve(os.path.join(src_dir, "ppt", "slideMasters"), target)
                    new_path = self._copy_media_file(src_path)
                    if new_path:
                        rel.set('Target', f"../{new_path}")

                elif rel_type == REL_THEME:
                    src_theme = _resolve(os.path.join(src_dir, "ppt", "slideMasters"), target)
                    if os.path.exists(src_theme):
                        new_theme = f"theme{self._next_theme}.xml"
                        self._next_theme += 1
                        dst_theme = os.path.join(self.base_dir, "ppt", "theme", new_theme)
                        os.makedirs(os.path.dirname(dst_theme), exist_ok=True)
                        shutil.copy2(src_theme, dst_theme)

                        src_theme_rels = os.path.join(
                            src_dir, "ppt", "theme", "_rels",
                            os.path.basename(src_theme) + ".rels")
                        if os.path.exists(src_theme_rels):
                            dst_theme_rels = os.path.join(
                                self.base_dir, "ppt", "theme", "_rels",
                                f"{new_theme}.rels")
                            theme_rels_tree = _parse_xml(src_theme_rels)
                            self._remap_rels_media(
                                theme_rels_tree.getroot(),
                                os.path.join(src_dir, "ppt", "theme"))
                            os.makedirs(os.path.dirname(dst_theme_rels), exist_ok=True)
                            _write_xml(theme_rels_tree, dst_theme_rels)

                        rel.set('Target', f"../theme/{new_theme}")
                        self._add_ct_override(f"/ppt/theme/{new_theme}", THEME_CT)

                elif rel_type == REL_LAYOUT:
                    layout_rid_map[rid] = os.path.basename(target)

            for rel in rels_root.findall(f'{{{NS_REL}}}Relationship'):
                rid = rel.get('Id', '')
                if rid in layout_rid_map:
                    new_layout = self._copy_layout_for_master(
                        source_pptx, src_dir, layout_rid_map[rid], new_master)
                    rel.set('Target', f"../slideLayouts/{new_layout}")

            os.makedirs(os.path.dirname(dst_rels_path), exist_ok=True)
            _write_xml(rels_tree, dst_rels_path)

        for li in master_xml_root.findall(f'.//{{{NS_P}}}sldLayoutId'):
            li.set('id', str(self._next_layout_id))
            self._next_layout_id += 1

        _write_xml(master_xml_tree, dst_master_path)

        self._add_ct_override(f"/ppt/slideMasters/{new_master}", MASTER_CT)

        rid = self._alloc_pres_rid()
        pres_rels_root = self.rels_tree.getroot()
        rel_el = etree.SubElement(pres_rels_root, f'{{{NS_REL}}}Relationship')
        rel_el.set('Id', rid)
        rel_el.set('Type', REL_MASTER)
        rel_el.set('Target', f"slideMasters/{new_master}")

        pres_root = self.pres_tree.getroot()
        masterIdLst = pres_root.find(f'{{{NS_P}}}sldMasterIdLst')
        if masterIdLst is None:
            sldIdLst = pres_root.find(f'{{{NS_P}}}sldIdLst')
            masterIdLst = etree.SubElement(pres_root, f'{{{NS_P}}}sldMasterIdLst')
            if sldIdLst is not None:
                pres_root.remove(masterIdLst)
                pres_root.insert(list(pres_root).index(sldIdLst), masterIdLst)

        master_id_el = etree.SubElement(masterIdLst, f'{{{NS_P}}}sldMasterId')
        master_id_el.set('id', str(self._next_layout_id))
        self._next_layout_id += 1
        master_id_el.set(f'{{{NS_R}}}id', rid)

        self._master_map[key] = new_master
        if self._verbose:
            print(f"    [master] {master_filename} -> {new_master} "
                  f"({len(layout_rid_map)} layouts)")
        return new_master

    def _copy_layout_for_master(self, source_pptx, src_dir, layout_filename, new_master):
        key = (str(source_pptx), layout_filename)
        if key in self._layout_map:
            return self._layout_map[key]

        new_layout = f"slideLayout{self._next_layout}.xml"
        self._next_layout += 1

        src_path = os.path.join(src_dir, "ppt", "slideLayouts", layout_filename)
        dst_path = os.path.join(self.base_dir, "ppt", "slideLayouts", new_layout)
        os.makedirs(os.path.dirname(dst_path), exist_ok=True)
        shutil.copy2(src_path, dst_path)

        src_rels = os.path.join(src_dir, "ppt", "slideLayouts", "_rels",
                                f"{layout_filename}.rels")
        dst_rels = os.path.join(self.base_dir, "ppt", "slideLayouts", "_rels",
                                f"{new_layout}.rels")

        if os.path.exists(src_rels):
            rels_tree = _parse_xml(src_rels)
            rels_root = rels_tree.getroot()

            for rel in rels_root.findall(f'{{{NS_REL}}}Relationship'):
                rel_type = rel.get('Type', '')
                target = rel.get('Target', '')

                if _is_media_rel(rel_type):
                    src_media = _resolve(
                        os.path.join(src_dir, "ppt", "slideLayouts"), target)
                    new_media = self._copy_media_file(src_media)
                    if new_media:
                        rel.set('Target', f"../{new_media}")

                elif rel_type == REL_MASTER:
                    rel.set('Target', f"../slideMasters/{new_master}")

            os.makedirs(os.path.dirname(dst_rels), exist_ok=True)
            _write_xml(rels_tree, dst_rels)

        self._add_ct_override(f"/ppt/slideLayouts/{new_layout}", LAYOUT_CT)
        self._layout_map[key] = new_layout
        return new_layout

    def _get_layout_for_slide(self, source_pptx, src_dir, slide_rels_root):
        for rel in slide_rels_root.findall(f'{{{NS_REL}}}Relationship'):
            if rel.get('Type') == REL_LAYOUT:
                layout_file = os.path.basename(rel.get('Target', ''))
                key = (str(source_pptx), layout_file)
                if key in self._layout_map:
                    return self._layout_map[key]

                layout_rels_path = os.path.join(
                    src_dir, "ppt", "slideLayouts", "_rels",
                    f"{layout_file}.rels")
                if os.path.exists(layout_rels_path):
                    layout_rels = _parse_xml(layout_rels_path)
                    for lrel in layout_rels.getroot().findall(
                            f'{{{NS_REL}}}Relationship'):
                        if lrel.get('Type') == REL_MASTER:
                            master_file = os.path.basename(lrel.get('Target', ''))
                            self._copy_master_with_all_layouts(
                                source_pptx, src_dir, master_file)
                            return self._layout_map.get(key)
        return None

    def add_slide(self, source_pptx: Path, slide_index: int) -> bool:
        """
        Copy a slide from source_pptx into the merged presentation.

        Args:
            source_pptx: Path to the source PPTX file.
            slide_index: Zero-based index of the slide to copy.

        Returns:
            True if the slide was copied successfully, False otherwise.
        """
        src_dir = self._extract_source(source_pptx)
        src_slide_file = f"slide{slide_index + 1}.xml"
        src_slide_path = os.path.join(src_dir, "ppt", "slides", src_slide_file)

        if not os.path.exists(src_slide_path):
            return False

        new_slide_file = f"slide{self._next_slide}.xml"
        self._next_slide += 1

        dst_slide = os.path.join(self.base_dir, "ppt", "slides", new_slide_file)
        os.makedirs(os.path.dirname(dst_slide), exist_ok=True)
        shutil.copy2(src_slide_path, dst_slide)

        src_rels = os.path.join(src_dir, "ppt", "slides", "_rels",
                                f"{src_slide_file}.rels")
        dst_rels = os.path.join(self.base_dir, "ppt", "slides", "_rels",
                                f"{new_slide_file}.rels")

        if os.path.exists(src_rels):
            rels_tree = _parse_xml(src_rels)
            rels_root = rels_tree.getroot()

            # Strip comment references (we don't copy comment files)
            for rel in list(rels_root.findall(f'{{{NS_REL}}}Relationship')):
                target = rel.get('Target', '')
                if 'comment' in target.lower() or 'comment' in rel.get('Type', '').lower():
                    rels_root.remove(rel)

            new_layout = self._get_layout_for_slide(
                source_pptx, src_dir, rels_root)

            for rel in rels_root.findall(f'{{{NS_REL}}}Relationship'):
                target = rel.get('Target', '')
                rel_type = rel.get('Type', '')

                if _is_media_rel(rel_type):
                    src_path = _resolve(
                        os.path.join(src_dir, "ppt", "slides"), target)
                    new_path = self._copy_media_file(src_path)
                    if new_path:
                        rel.set('Target', f"../{new_path}")

                elif rel_type == REL_LAYOUT and new_layout:
                    rel.set('Target', f"../slideLayouts/{new_layout}")

                elif rel_type == REL_NOTES:
                    self._copy_notes_master(src_dir)
                    src_notes_file = os.path.basename(target)
                    src_notes = os.path.join(
                        src_dir, "ppt", "notesSlides", src_notes_file)
                    if os.path.exists(src_notes):
                        new_notes_file = f"notesSlide{self._next_notes}.xml"
                        self._next_notes += 1
                        dst_notes = os.path.join(
                            self.base_dir, "ppt", "notesSlides", new_notes_file)
                        os.makedirs(os.path.dirname(dst_notes), exist_ok=True)
                        shutil.copy2(src_notes, dst_notes)
                        rel.set('Target', f"../notesSlides/{new_notes_file}")

                        src_nrels = os.path.join(
                            src_dir, "ppt", "notesSlides", "_rels",
                            f"{src_notes_file}.rels")
                        if os.path.exists(src_nrels):
                            nrels_tree = _parse_xml(src_nrels)
                            for nrel in nrels_tree.getroot().findall(
                                    f'{{{NS_REL}}}Relationship'):
                                ntype = nrel.get('Type', '')
                                if 'slide' in ntype and 'notesSlide' not in ntype and 'notesMaster' not in ntype:
                                    nrel.set('Target', f"../slides/{new_slide_file}")
                                elif _is_media_rel(ntype):
                                    sp = _resolve(
                                        os.path.join(src_dir, "ppt", "notesSlides"),
                                        nrel.get('Target', ''))
                                    np = self._copy_media_file(sp)
                                    if np:
                                        nrel.set('Target', f"../{np}")
                            dst_nrels = os.path.join(
                                self.base_dir, "ppt", "notesSlides", "_rels",
                                f"{new_notes_file}.rels")
                            os.makedirs(os.path.dirname(dst_nrels), exist_ok=True)
                            _write_xml(nrels_tree, dst_nrels)

                        self._add_ct_override(
                            f"/ppt/notesSlides/{new_notes_file}", NOTES_CT)

            os.makedirs(os.path.dirname(dst_rels), exist_ok=True)
            _write_xml(rels_tree, dst_rels)

        pres_root = self.pres_tree.getroot()
        sldIdLst = pres_root.find(f'{{{NS_P}}}sldIdLst')
        if sldIdLst is None:
            sldIdLst = etree.SubElement(pres_root, f'{{{NS_P}}}sldIdLst')

        rid = self._alloc_pres_rid()
        sldId = etree.SubElement(sldIdLst, f'{{{NS_P}}}sldId')
        sldId.set('id', str(self._next_slide_id))
        sldId.set(f'{{{NS_R}}}id', rid)
        self._next_slide_id += 1

        pres_rels_root = self.rels_tree.getroot()
        rel_el = etree.SubElement(pres_rels_root, f'{{{NS_REL}}}Relationship')
        rel_el.set('Id', rid)
        rel_el.set('Type', REL_SLIDE)
        rel_el.set('Target', f"slides/{new_slide_file}")

        self._add_ct_override(f"/ppt/slides/{new_slide_file}", SLIDE_CT)
        self._slide_count += 1
        return True

    def _fix_presentation_element_order(self):
        pres_root = self.pres_tree.getroot()
        order_map = {name: i for i, name in enumerate(PRES_ELEMENT_ORDER)}

        children = list(pres_root)
        for child in children:
            pres_root.remove(child)

        for child in sorted(children, key=lambda el: order_map.get(
                etree.QName(el).localname, 999)):
            pres_root.append(child)

    def save(self, output_path: Path):
        """Write the merged presentation to output_path."""
        self._fix_presentation_element_order()
        _write_xml(self.pres_tree, self.pres_xml_path)
        _write_xml(self.rels_tree, self.pres_rels_path)
        _write_xml(self.ct_tree, self.ct_path)

        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for root, dirs, files in os.walk(self.base_dir):
                for f in files:
                    fp = os.path.join(root, f)
                    arcname = os.path.relpath(fp, self.base_dir)
                    zf.write(fp, arcname)

        if self._verbose:
            size_kb = os.path.getsize(output_path) / 1024
            print(f"Saved: {output_path.name} ({size_kb:.0f} KB, "
                  f"{self._slide_count} slides)")

    def cleanup(self):
        """Remove temporary files."""
        shutil.rmtree(self.tmp_dir, ignore_errors=True)

    def __enter__(self):
        return self

    def __exit__(self, *args):
        self.cleanup()
